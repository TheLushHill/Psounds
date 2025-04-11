# 该文件夹中定义前端向后端发送文件/请求，后端判断是哪种文件类型并调用该文件的api

import glob
import logging, warnings, tempfile
from modelscope.utils.logger import get_logger
logging.basicConfig(level=logging.WARNING)
modelscope_logger = get_logger()
modelscope_logger.setLevel(logging.ERROR)  # 设置为 ERROR 或更高
modelscope_logger.propagate = False

from flask import Flask, request, jsonify, render_template, session, Response
import os, sys, requests, json, io
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(os.path.join(os.path.dirname(__file__), "Model"))
sys.path.append(os.path.join(os.path.dirname(__file__), "Model/Training/GS_Model"))


import argparse, shutil
from multiprocessing import freeze_support
from pathlib import Path

# from Model.tts_main import get_audio
# import soundfile as sf

app = Flask(__name__)

ALLOWED_EXTENSIONS = {'docx', 'pptx'}
ALLOWED_MIME_TYPES = {  # 这些是HTTPmime类型,特定文件的代号
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
}

# 适配中文的扩展名！
def allowed_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


@app.route('/upload', methods=['POST'])                             # JS的AJAX请求，fetch("/upload", {...})
def upload_file():
    from GetWord.docx.DocxWord import get_docxtext
    from GetWord.pptx.PPTWord import get_ppttext

    file = request.files['file']
    filename = file.filename

    # 检查文件名是否为空
    if file.filename == '':
        return jsonify({"error": "空文件名"}), 400

    # 校验扩展名
    if not allowed_file(filename):
        return jsonify({"error": "文件扩展名不合法"}), 400

    # 删除文件内容检测部分，直接通过扩展名匹配
    file_extension = filename.rsplit('.', 1)[1].lower()
    detected_type = ALLOWED_MIME_TYPES.get(file_extension)

    # 如果扩展名不在允许列表中，直接拦截
    if not detected_type:
        return jsonify({"error": "文件类型不合法"}), 400
    header = file.stream.read(4)
    file.seek(0)   # 重置文件指针
    if header != b'PK\x03\x04':
        return "Invalid PPTX (not a ZIP file)", 400
    
    # 如果是PPT文件，则保存
    # if detected_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
    #     file.save(os.path.join("PPT", filename))

    # 6. 如果文件过大(大于1G)，则保存在本地文件
    # 先创建缓存文件夹以保存文件
    # if( >= 1000*1024):
    #   path1 = os.getcwd()
    #   os.makedirs(".Cache\Input", exist_ok=True)
    #   os.makedirs(".Cache\Output", exist_ok=True)
    #   UPLOAD_FOLDER = os.path.join(path1, ".Cache\Input")
    #   POST_FOLDER = os.path.join(path1, ".Cache\Output")
    #
    # 判断文件夹生成是否成功
    #   if not os.path.exists(UPLOAD_FOLDER):
    #       os.makedirs(UPLOAD_FOLDER)
    # 在本地保存
    #   file.save(os.path.join(UPLOAD_FOLDER, filename))


# 在这里判断文件类型，并调用各类型的api   后续在这里判断一下内存大小
    if detected_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return jsonify(get_docxtext(file))
    elif detected_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return jsonify(get_ppttext(file))


# 解析成功返回前端。前端发送一段文字后，在这里进行预处理

# 前端发送过来的人物音频，先保存到Model\\input_audio
@app.route('/SaveAudio', methods=['POST'])
def SaveAudio():

    ALLOWED_AUDIO_EXTENSIONS = {'wav', 'mp3', 'm4a'}
    def allowed_audio_file(filename):
        return '.' in filename and \
            filename.rsplit('.', 1)[1].lower() in ALLOWED_AUDIO_EXTENSIONS
    
    save_dir = os.path.abspath("Model\\input_audio")
    audio_files = request.files.getlist('file')  # 获取多个文件

    for file in audio_files:
        if file:
            filename = file.filename
            file_path = os.path.join(save_dir, filename)
            if allowed_audio_file(file.filename):   
                file.save(file_path)
                print(f"音频文件 {filename} 保存成功")
            else:
                return jsonify({"error": "音频文件类型不合法"}), 400
        else:
            return jsonify({"error": "没有找到文件或文件名空"}), 400
           
    return jsonify({"状态": "音频保存成功" , "文件夹路径": "Model\\input_audio"}), 200
    
# 解析成功,把前端刚刚上传的文件名依次返回前端
@app.route('/getfilename', methods=['POST'])
def get_filename():
    try:
        # 获取目录下所有文件名
        files = os.listdir("Model\\input_audio")
        # 过滤出文件（排除目录）
        files = [f for f in files if os.path.isfile(os.path.join("Model\\input_audio", f))]
        return jsonify({'filenames': files}), 200
    except FileNotFoundError:
        return jsonify({'error': 'Directory not found'}), 404


##数据集的预处理
@app.route('/preprocess', methods=['POST'])
def preprocess():
    from Model.Tools.uvr5.uvrmain import uvr
    from Model.Tools.slice_audio.slice_audio import slice_audio
    from Model.Tools.cmd_denoise.cmd_denoise import execute_denoise
    from Model.Tools.asr.funasr_asr import execute_asr

    result1 = uvr("Model\\input_audio",
                 "Model\\Tools\\output_cache\\uvr5_opt\\vocal",
                 "Model\\Tools\\output_cache\\uvr5_opt\\instrumental"
                 )
    print(result1)

    result2 = slice_audio("Model\\Tools\\output_cache\\uvr5_opt\\vocal", "Model\\Tools\\output_cache\\slicer_opt")
    print(result2)

    result3 = execute_denoise("Model\\Tools\\output_cache\\slicer_opt", "Model\\Tools\\output_cache\\denoise_opt")
    print(result3)

    result4 = execute_asr("Model\\Tools\\output_cache\\denoise_opt", "Model\\Tools\\output_cache\\asr_opt")
    print(result4)

    if result4:
        return jsonify({"log": "预处理完成"}), 200
    else:
        return jsonify({"error": "预处理失败"}), 500

##之后要重定向到打标界面，手动进行修正
# @app.route('/subfix', methods=['POST'])
# def subfix_editor():
#     subfix.start()
#     return jsonify({"log": "打标界面已启动"}), 200

# 格式化音频
@app.route('/format', methods=['POST'])
def format_audio():
    from Model.Training.prepare_datasets.gettext1 import slice2bert
    from Model.Training.prepare_datasets.gethubert2 import get_hubert
    from Model.Training.prepare_datasets.getsemanic3 import get_semanic

    character = request.json.get("character")
    slice2bert("Model/Tools/output_cache/asr_opt/denoise_opt.list", character)
    get_hubert(character)
    get_semanic(character)
    
    return jsonify({"log": "音频格式化成功，可以训练模型"}), 200



# 训练模型
@app.route('/train', methods=['POST'])
def train():
    from Model.Training.GS_Model.s1_train import main as s1_main
    from Model.Training.GS_Model.s2_train import main as s2_main

    #调用 s1_train 的 main 函数的参数
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "-c",
        "--config_file",
        type=str,
        default="Model/Training/GS_Model/configs/s1longer_v2.yaml",
        help="path of config file",
    )
    args = parser.parse_args()
    character = request.json.get("character") 

    s1_main(args, character)
    s2_main(character)

    # 移动.ckpt与.pth文件至父级文件夹，并删除其他不要的文件夹
    dest_dir = Path(f"Model/trained/{character}")
    ckpt_dir = Path(f"Model/trained/{character}/ckpt")
    eval_dir = Path(f"Model/trained/{character}/eval")

    target_files = glob.glob(str(ckpt_dir / "*epoch=11*.ckpt"))
    src_path = Path(target_files[0])  # 取首个匹配项
    dest_path = dest_dir / src_path.name

    shutil.move(str(src_path),str(dest_path))

    shutil.rmtree(f"Model\\trained\\{character}\\{character}")
    shutil.rmtree(ckpt_dir)
    shutil.rmtree(eval_dir)

    s1delete = glob.glob(str(dest_dir /"*e4*.pth"))
    d1 = Path(s1delete[0])
    s2delete = glob.glob(str(dest_dir /"*e8*.pth"))
    d2 = Path(s2delete[0])
    os.remove(d1)
    os.remove(d2)
    
    for file_path in glob.glob("*17*.pth"):
        if os.path.isfile(file_path):
            os.remove(file_path)
            print(f"已删除pth缓存文件: {file_path}")

    if os.path.exists(f"Model\\trained\\{character}"):
        return jsonify({"log": "模型训练成功"}), 200
    else:
        return jsonify({"log": "未找到模型路径"}), 500


# 训练完成后，在GPT_weights、SoVITS_weights中提取模型，进行TTS
#tts 开始
# from Model.tts_main import get_audio
# import soundfile as sf
def process_audio_sf(sr, audio_data):
    with io.BytesIO() as buffer:
        sf.write(buffer, audio_data, sr, format='WAV')
        return buffer.getvalue()

# 获取处理后的音频数据
def get_processed_audio(data, streaming=False):
    """获取处理后的音频数据"""
    
    raw_audio = get_audio(data, streaming=streaming)
    
    if not streaming:
        for sr, audio_data in raw_audio:
            yield process_audio_sf(sr, audio_data)
    else:
        # 处理流式音频
        for chunk in raw_audio:
            yield chunk

# 返回整个 wav 文件
@app.route("/tts", methods=["POST"])
def handletts():

    data = request.json;  
    print("路徑：")
    print(sys.path)
    data = {
        "text": data.get("text", "你好，世界"),
        "character": data.get("character", "Hutao"),
        "prompt_text": data.get("prompt_text", "我说白术，你不会看不出来吧？难不成你师父，忘了教你这门功夫"),
        "ref_audio_path": data.get("ref_audio_path", "Model/trained/Hutao/我说白术，你不会看不出来吧？难不成你师父，忘了教你这门功夫？.wav"),
    }
    
    def generate():
        for chunk in get_processed_audio(data):
            yield chunk;

    return Response(generate(), status=200,mimetype="audio/wav");

# 流式返回
@app.route("/tts_stream", methods=["POST"])
def handletts_stream():

    data = request.json;
    
    data = {
        "text": data.get("text", "你好，世界"),
        "character": data.get("character", "Hutao"),
        "prompt_text": data.get("prompt_text", "我说白术，你不会看不出来吧？难不成你师父，忘了教你这门功夫"),
        "ref_audio_path": data.get("ref_audio_path", "Model/trained/Hutao/我说白术，你不会看不出来吧？难不成你师父，忘了教你这门功夫？.wav"),
    }

    def generate():
        for chunk in get_processed_audio(data, streaming=True):
            print(f"Yielding chunk of size: {len(chunk)}")  # 打印每个块的大小
            yield chunk;

    return Response(generate(), status=200,mimetype="audio/wav");
# tts 结束


#当用户提供的是PPT文件时，在每一页中添加音频
@app.route('/PPTaudio', methods=['POST'])
def PPTaudio():
    from GetWord.pptx.PPTWord import get_ppttext
    from pptx import Presentation
    from pptx.util import Inches

    file = request.files['file']
    filename = file.filename
    file_extension = filename.rsplit('.', 1)[1].lower()
    detected_type = ALLOWED_MIME_TYPES.get(file_extension)

    file.stream.seek(0)
    if detected_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        slide_texts = get_ppttext(file)
    else:
        return jsonify({"error": "文件类型不合法"}), 400

    character = request.form.get("character")

    file.stream.seek(0)
    try:
        prs = Presentation(file)
    except Exception as e:
        return jsonify({"error": f"PPT文件读取失败: {str(e)}"}), 500

    # 遍历每一页，生成并插入音频
    for idx, slide in enumerate(prs.slides):
        text = "\n".join(slide_texts[idx]).strip()
        if not text:
            continue

        # 3) 收集 TTS 输出的所有块，拼成 bytes
        try:
            chunks = list(get_processed_audio({"text": text, "character": character}))
            audio_bytes = b"".join(chunks)
        except Exception as e:
            print(f"TTS 失败: {e}")
            continue

        # 4) 写临时 mp3 并插入
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            tmp.write(audio_bytes)
            tmp.flush()
            tmp_path = tmp.name

        try:
            audio_shape = slide.shapes.add_movie(
                tmp_path,
                Inches(1), Inches(1),
                Inches(1), Inches(1),
                mime_type='audio/wav',
                poster_frame_image=None
            )
            audio_shape.media_format.show_controls = False
        except Exception as e:
            print(f"插入失败: {e}")
        finally:
            os.remove(tmp_path)

    # 5) 保存并返回
    os.makedirs('PPT', exist_ok=True)
    out_path = os.path.join('PPT', 'modified_' + file.filename)
    prs.save(out_path)
    return jsonify({"output_path": out_path}), 200


if __name__ == '__main__':
    app.run(host="127.0.0.1", port=8000, debug=True, use_reloader=False)
